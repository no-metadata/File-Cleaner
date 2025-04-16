import sys
import os
import shutil
import datetime
import zipfile
import xml.etree.ElementTree as ET
import json
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
from pydub import AudioSegment
import pikepdf
from lxml import etree

dt = datetime.datetime(1900, 1, 1)

WORD_DEFAULTS = {
    'author': "",
    'category': "",
    'comments': "",
    'content_status': "",
    'created': dt,
    'identifier': "",
    'keywords': "",
    'language': "",
    'last_modified_by': "",
    'last_printed': dt,
    'modified': dt,
    'revision': 1,
    'subject': "",
    'title': "",
    'version': "",
}

EXCEL_DEFAULTS = {
    'creator': "",
    'title': "",
    'description': "",
    'subject': "",
    'keywords': "",
    'category': "",
    'lastModifiedBy': "",
    'created': dt,
    'modified': dt,
    'lastPrinted': dt,
    'revision': 1,
    'manager': "",
    'company': "",
}

def clean_ooxml_extended_properties(file_path):
    ext = get_file_type(file_path)
    if ext not in ['docx', 'xlsx', 'pptx']:
        return

    if ext == 'docx':
        props_to_clear = {
            'Company': "",
            'Manager': "",
            'TotalTime': "0",
            'AppVersion': "1.0",
        }
    elif ext == 'xlsx':
        props_to_clear = {
            'Company': "",
            'Manager': "",
            'AppVersion': "1.0",
        }
    elif ext == 'pptx':
        props_to_clear = {
            'Company': "",
            'Manager': "",
            'AppVersion': "1.0",
        }
    try:
        temp_file = file_path + '.tmp'
        ns = {'ep': 'http://schemas.openxmlformats.org/officeDocument/2006/extended-properties'}
        with zipfile.ZipFile(file_path, 'r') as zin:
            with zipfile.ZipFile(temp_file, 'w') as zout:
                for item in zin.infolist():
                    if item.filename == 'docProps/app.xml':
                        xml_data = zin.read(item.filename)
                        root = ET.fromstring(xml_data)
                        for tag, default_value in props_to_clear.items():
                            elem = root.find('ep:' + tag, ns)
                            if elem is not None:
                                elem.text = default_value
                        new_xml_data = ET.tostring(root, encoding="utf-8", method="xml")
                        zout.writestr(item, new_xml_data)
                    else:
                        zout.writestr(item, zin.read(item.filename))
        os.replace(temp_file, file_path)
    except Exception as e:
        raise ValueError(f"Error cleaning extended properties: {e}")

def get_file_type(file_path):
    try:
        with open(file_path, "rb") as f:
            header = f.read(12)
    except Exception:
        return "unknown"

    # PDF
    if header.startswith(b"%PDF"):
        return "pdf"
    # PNG
    if header.startswith(b"\x89PNG"):
        return "png"
    # JPEG
    if header.startswith(b"\xff\xd8"):
        return "jpg"
    # FLAC
    if header.startswith(b"fLaC"):
        return "flac"
    # WAV
    if header.startswith(b"RIFF"):
        if header[8:12] == b"WAVE":
            return "wav"
    # MP3
    if header.startswith(b"ID3") or header[0:2] == b'\xff\xfb':
        return "mp3"
    # MS Office
    if header.startswith(b"PK\x03\x04"):
        try:
            with zipfile.ZipFile(file_path, 'r') as z:
                names = z.namelist()
                if "[Content_Types].xml" in names:
                    if any(name.startswith("word/") for name in names):
                        return "docx"
                    elif any(name.startswith("xl/") for name in names):
                        return "xlsx"
                    elif any(name.startswith("ppt/") for name in names):
                        return "pptx"
        except Exception:
            return "unknown"
    return "unknown"

def clean_word_document(file_path):
    doc = Document(file_path)
    for prop, default_value in WORD_DEFAULTS.items():
        if hasattr(doc.core_properties, prop):
            setattr(doc.core_properties, prop, default_value)
    document_element = doc.element
    for comment_ref in document_element.xpath('//w:commentReference | //w:commentRangeStart | //w:commentRangeEnd'):
        comment_ref.getparent().remove(comment_ref)
    if 'comments' in doc.part.rels:
        del doc.part.rels['comments']
    doc.save(file_path)
    clean_ooxml_extended_properties(file_path)

def clean_excel_document(file_path):
    wb = load_workbook(file_path)
    properties = wb.properties
    for prop, default_value in EXCEL_DEFAULTS.items():
        if hasattr(properties, prop):
            setattr(properties, prop, default_value)
    for sheet in wb:
        for row in sheet:
            for cell in row:
                cell.comment = None
    wb.save(file_path)
    clean_ooxml_extended_properties(file_path)

def clean_powerpoint_document(file_path):
    prs = Presentation(file_path)
    for prop, default_value in WORD_DEFAULTS.items():
        if hasattr(prs.core_properties, prop):
            setattr(prs.core_properties, prop, default_value)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type in [MSO_SHAPE_TYPE.COMMENT, MSO_SHAPE_TYPE.INK_COMMENT]:
                sp = shape._element
                sp.getparent().remove(sp)
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = ''
    prs.save(file_path)
    clean_ooxml_extended_properties(file_path)

def clean_image(file_path):
    try:
        img = Image.open(file_path)
        data = list(img.getdata())
        image_no_metadata = Image.new(img.mode, img.size)
        image_no_metadata.putdata(data)
        image_no_metadata.save(file_path)
    except Exception as e:
        raise ValueError(f"Error cleaning image: {e}")

def clean_audio_file(file_path):
    try:
        audio = AudioSegment.from_file(file_path)
        ext = os.path.splitext(file_path)[1][1:]
        audio.export(file_path, format=ext)
    except Exception as e:
        raise ValueError(f"Error cleaning audio file: {e}")

def clean_pdf_file(file_path):
    try:
        with pikepdf.open(file_path, allow_overwriting_input=True) as pdf:
            for key in list(pdf.docinfo.keys()):
                del pdf.docinfo[key]
            if '/Metadata' in pdf.Root:
                del pdf.Root['/Metadata']
            pdf.save(file_path)
    except Exception as e:
        raise ValueError(f"Error cleaning PDF file: {e}")

def clean_file(file_path):
    backup = file_path + ".bak"
    try:
        shutil.copy2(file_path, backup)
        file_type = get_file_type(file_path)
        if file_type == 'docx':
            clean_word_document(file_path)
        elif file_type == 'xlsx':
            clean_excel_document(file_path)
        elif file_type == 'pptx':
            clean_powerpoint_document(file_path)
        elif file_type in ['png', 'jpg']:
            clean_image(file_path)
        elif file_type in ['wav', 'mp3', 'flac']:
            clean_audio_file(file_path)
        elif file_type == 'pdf':
            clean_pdf_file(file_path)
        else:
            raise ValueError("Unsupported file type")
        os.remove(backup)
    except Exception as e:
        try:
            shutil.copy2(backup, file_path)
            os.remove(backup)
        except Exception as revert_err:
            raise ValueError(f"Failed to revert {file_path} after error: {revert_err}") from e
        raise ValueError(f"Error cleaning {file_path}: {e}") from e

SETTINGS_FILE = os.path.join(os.path.expanduser("~"), ".file_cleaner_settings.json")

def load_settings():
    default_settings = {"language": "en"}
    if os.path.exists(SETTINGS_FILE):
        try:
            with open(SETTINGS_FILE, "r", encoding="utf-8") as f:
                settings = json.load(f)
            return settings
        except Exception:
            return default_settings
    else:
        return default_settings

def save_settings(settings):
    try:
        with open(SETTINGS_FILE, "w", encoding="utf-8") as f:
            json.dump(settings, f, indent=4)
    except Exception as e:
        print(f"Warning: could not save settings: {e}")

from PyQt5.QtWidgets import QListWidget, QAbstractItemView

class FileListWidget(QListWidget):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setAcceptDrops(True)
        self.setDragEnabled(True)
        self.setDropIndicatorShown(True)
        self.setDragDropMode(QAbstractItemView.DragDrop)  # Enable drag and drop mode

    def dragEnterEvent(self, event):
        if event.mimeData().hasUrls():
            event.acceptProposedAction()
        else:
            event.ignore()

    def dragMoveEvent(self, event):
        if event.mimeData().hasUrls():
            event.acceptProposedAction()
        else:
            event.ignore()

    def dropEvent(self, event):
        if event.mimeData().hasUrls():
            for url in event.mimeData().urls():
                file_path = url.toLocalFile()
                if file_path and file_path not in self.get_file_list():
                    self.addItem(file_path)
            event.acceptProposedAction()

    def get_file_list(self):
        return [self.item(i).text() for i in range(self.count())]

from PyQt5.QtWidgets import (
    QApplication, QMainWindow, QPushButton, QVBoxLayout, QWidget,
    QFileDialog, QMessageBox, QHBoxLayout, QProgressBar, QComboBox, QLabel
)
from PyQt5.QtCore import Qt

class FileCleanerUI(QMainWindow):
    def __init__(self):
        super().__init__()
        self.settings = load_settings()
        self.current_lang = self.settings.get("language", "en")
        self.language_strings = {
            "en": {
                "window_title": "File Cleaner",
                "add_files": "Add Files",
                "clean_files": "Clean Files",
                "clear_list": "Clear List",
                "no_files": "No Files",
                "please_add_files": "Please add files to clean.",
                "cleaning_completed": "Cleaning completed.\n\nSuccessful: {0}\nFailed: {1}",
                "failed_files": "Failed files:\n",
                "user_guide": "User Guide:\n\nThis application cleans metadata from Microsoft Office documents, images, audio, and PDF files.\nDrag and drop files or click 'Add Files' to begin.\nUse the language switch to change the interface language.",
                "user_guide_button": "User Guide"
            },
            "ru": {
                "window_title": "Очиститель файлов",
                "add_files": "Добавить файлы",
                "clean_files": "Очистить файлы",
                "clear_list": "Очистить список",
                "no_files": "Нет файлов",
                "please_add_files": "Пожалуйста, добавьте файлы для очистки.",
                "cleaning_completed": "Очистка завершена.\n\nУспешно: {0}\nНе удалось: {1}",
                "failed_files": "Файлы с ошибками:\n",
                "user_guide": "Руководство пользователя:\n\nЭто приложение очищает метаданные из документов Microsoft Office, изображений, аудио и PDF файлов.\nПеретащите файлы или нажмите «Добавить файлы», чтобы начать.\nИспользуйте переключатель языка для смены языка интерфейса.",
                "user_guide_button": "Руководство пользователя"
            }
        }
        self.initUI()
        self.updateLanguage()

    def initUI(self):
        self.setGeometry(100, 100, 800, 600)
        central_widget = QWidget()
        self.setCentralWidget(central_widget)
        main_layout = QVBoxLayout()

        # Language switch
        lang_layout = QHBoxLayout()
        lang_label = QLabel("Language:")
        self.languageCombo = QComboBox()
        self.languageCombo.addItem("🇬🇧 English", "en")
        self.languageCombo.addItem("🇷🇺 Русский", "ru")
        index = self.languageCombo.findData(self.current_lang)
        if index >= 0:
            self.languageCombo.setCurrentIndex(index)
        self.languageCombo.currentIndexChanged.connect(self.onLanguageChanged)
        lang_layout.addWidget(lang_label)
        lang_layout.addWidget(self.languageCombo)
        lang_layout.addStretch()
        main_layout.addLayout(lang_layout)

        self.list_widget = FileListWidget()
        main_layout.addWidget(self.list_widget)

        # Progress bar
        self.progressBar = QProgressBar()
        self.progressBar.setValue(0)
        main_layout.addWidget(self.progressBar)

        # Buttons
        button_layout = QHBoxLayout()
        self.add_button = QPushButton()
        self.add_button.clicked.connect(self.add_files)
        button_layout.addWidget(self.add_button)
        self.clean_button = QPushButton()
        self.clean_button.clicked.connect(self.clean_files)
        button_layout.addWidget(self.clean_button)
        self.clear_button = QPushButton()
        self.clear_button.clicked.connect(self.clear_list)
        button_layout.addWidget(self.clear_button)
        # User Guide Button
        self.guide_button = QPushButton()
        self.guide_button.clicked.connect(self.show_user_guide)
        button_layout.addWidget(self.guide_button)
        main_layout.addLayout(button_layout)

        central_widget.setLayout(main_layout)

    def onLanguageChanged(self):
        self.current_lang = self.languageCombo.currentData()
        self.updateLanguage()
        self.settings["language"] = self.current_lang
        save_settings(self.settings)

    def updateLanguage(self):
        strings = self.language_strings[self.current_lang]
        self.setWindowTitle(strings["window_title"])
        self.add_button.setText(strings["add_files"])
        self.clean_button.setText(strings["clean_files"])
        self.clear_button.setText(strings["clear_list"])
        self.guide_button.setText(strings["user_guide_button"])

    def add_files(self):
        files, _ = QFileDialog.getOpenFileNames(self, "Select Files")
        if files:
            for file_path in files:
                if file_path not in self.list_widget.get_file_list():
                    self.list_widget.addItem(file_path)

    def clear_list(self):
        self.list_widget.clear()

    def show_user_guide(self):
        strings = self.language_strings[self.current_lang]
        QMessageBox.information(self, strings["user_guide_button"], strings["user_guide"])

    def clean_files(self):
        files = self.list_widget.get_file_list()
        strings = self.language_strings[self.current_lang]
        if not files:
            QMessageBox.warning(self, strings["no_files"], strings["please_add_files"])
            return
        success_files = []
        failed_files = []
        self.progressBar.setMaximum(len(files))
        self.progressBar.setValue(0)
        for i, file_path in enumerate(files, start=1):
            try:
                clean_file(file_path)
                success_files.append(file_path)
            except Exception as e:
                failed_files.append((file_path, str(e)))
            self.progressBar.setValue(i)
            QApplication.processEvents()
        msg = strings["cleaning_completed"].format(len(success_files), len(failed_files))
        if failed_files:
            msg += "\n" + strings["failed_files"] + "\n".join(f"{f}: {err}" for f, err in failed_files)
        self.clear_list()
        self.progressBar.setValue(0)
        QMessageBox.information(self, strings["window_title"], msg)

if __name__ == "__main__":
    app = QApplication(sys.argv)
    window = FileCleanerUI()
    window.show()
    sys.exit(app.exec_())
